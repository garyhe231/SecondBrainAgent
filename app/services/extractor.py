"""
Text extraction from various file types.
"""
import os
import subprocess
import tempfile
from pathlib import Path
from typing import Optional

SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm",
    ".pdf",
    ".docx",
    ".pptx",
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp",
    ".mp4", ".mov", ".avi", ".mkv", ".webm", ".m4v", ".mp3", ".wav", ".m4a",
}


def extract_text(file_path: str) -> Optional[str]:
    """Extract text from a file based on its extension."""
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        return None

    try:
        if ext in {".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm"}:
            return _extract_plaintext(file_path)
        elif ext == ".pdf":
            return _extract_pdf(file_path)
        elif ext == ".docx":
            return _extract_docx(file_path)
        elif ext == ".pptx":
            return _extract_pptx(file_path)
        elif ext in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}:
            return _extract_image(file_path)
        elif ext in {".mp4", ".mov", ".avi", ".mkv", ".webm", ".m4v", ".mp3", ".wav", ".m4a"}:
            return _extract_video_audio(file_path)
    except Exception as e:
        return f"[Extraction error: {e}]"

    return None


def _extract_plaintext(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _extract_pdf(file_path: str) -> str:
    import pdfplumber
    pages = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                pages.append(f"[Page {i+1}]\n{text}")
    return "\n\n".join(pages)


def _extract_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return "\n".join(paragraphs)


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_image(file_path: str) -> str:
    import pytesseract
    from PIL import Image
    img = Image.open(file_path)
    text = pytesseract.image_to_string(img)
    return text.strip() if text.strip() else "[No text found in image]"


def _extract_video_audio(file_path: str) -> str:
    import whisper
    ext = Path(file_path).suffix.lower()
    audio_path = file_path

    # If video, extract audio track first
    if ext in {".mp4", ".mov", ".avi", ".mkv", ".webm", ".m4v"}:
        tmp = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
        tmp.close()
        audio_path = tmp.name
        subprocess.run(
            ["ffmpeg", "-y", "-i", file_path, "-ac", "1", "-ar", "16000", audio_path],
            capture_output=True, check=True
        )

    model = whisper.load_model("base")
    result = model.transcribe(audio_path)

    if audio_path != file_path:
        os.unlink(audio_path)

    return result.get("text", "").strip()
